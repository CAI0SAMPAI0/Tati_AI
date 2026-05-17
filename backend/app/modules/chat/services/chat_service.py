"""
services/chat_service.py
Serviço para gerenciamento de chat, streaming de LLM e tarefas de background.
"""

import asyncio
import base64
import re
import io
import docx
import pypdf
from typing import Dict, List, Any, Optional
from datetime import date

from fastapi import WebSocket
from fastapi.concurrency import run_in_threadpool

from app.core.config import settings
from app.modules.chat.services.llm import stream_llm, text_to_speech, transcribe_audio
from app.shared.services.history import save_message, load_history, load_llm_history, auto_title
from app.modules.chat.services.prompt_builder import UserProfile, build_effective_prompt
from app.modules.chat.services.rag_search import obter_contexto_rag, RAGResult
from app.modules.activities.services.pronunciation_matcher import match_pronunciation
from app.core.dependencies.db import get_db
from fastapi import Depends
from supabase import Client


class ChatService:
    def __init__(self, db: Any = Depends(get_db)) -> None:
        if db is None or str(type(db)).find('Depends') != -1:
            from app.core.database import get_client
            self.db = get_client()
        else:
            self.db = db

    async def _execute_db(self, func, retries=3):
        """Helper para executar chamadas de banco com retry em caso de desconexão."""
        for attempt in range(retries):
            try:
                return await run_in_threadpool(func)
            except Exception as e:
                err_str = str(e).lower()
                if ('disconnected' in err_str or 'connection' in err_str or 'protocol' in err_str) and attempt < retries - 1:
                    print(f'[ChatService] DB connection issue, retrying ({attempt+1}/{retries})...')
                    await asyncio.sleep(0.5 * (attempt + 1))
                    continue
                raise e

    async def get_user_profile(self, username: str) -> UserProfile:
        def _fetch():
            return (
                self.db.table('users')
                .select('name, custom_prompt, level, focus')
                .eq('username', username)
                .limit(1)
                .execute()
                .data
            )

        try:
            rows = await self._execute_db(_fetch)
        except Exception as e:
            print(f'[ChatService] Error fetching user profile: {e}')
            rows = []
        
        data = rows[0] if rows else {}
        return UserProfile(
            username=username,
            name=data.get('name') or username,
            level=data.get('level') or 'Intermediate',
            focus=data.get('focus') or 'General Conversation',
            custom_prompt=(data.get('custom_prompt') or '').strip(),
        )

    async def check_access(self, username: str) -> Dict[str, Any]:
        """Verifica se o usuário pode enviar mensagens."""
        today = date.today()
        PAID_START = date(2026, 6, 30)
        FREE_MSG_LIMIT = 100

        def _fetch_user():
            return (
                self.db.table('users')
                .select(
                    'username, role, is_exempt, is_premium_active, created_at, free_messages_used'
                )
                .eq('username', username)
                .limit(1)
                .execute()
                .data
            )

        try:
            user_rows = await self._execute_db(_fetch_user)
        except Exception as e:
            print(f'[ChatService] Error checking access: {e}')
            return {'allowed': True, 'reason': None, 'free_messages_remaining': None}
        
        user = user_rows[0] if user_rows else {}

        is_admin = user.get('role') in settings.staff_roles
        is_exempt = user.get('is_exempt', False)

        if is_admin or is_exempt or user.get('is_premium_active'):
            return {'allowed': True, 'reason': None, 'free_messages_remaining': None}

        if today < PAID_START:
            user_created = date.fromisoformat(user['created_at'][:10])
            if user_created < PAID_START:
                return {
                    'allowed': True,
                    'reason': None,
                    'free_messages_remaining': None,
                }

        used = user.get('free_messages_used') or 0
        remaining = max(0, FREE_MSG_LIMIT - used)

        if remaining <= 0:
            return {
                'allowed': False,
                'reason': 'free_limit_reached',
                'free_messages_remaining': 0,
            }

        return {'allowed': True, 'reason': None, 'free_messages_remaining': remaining}

    async def process_chat_message(
        self,
        websocket: WebSocket,
        msg: Dict[str, Any],
        username: str,
        pending_drill_target: Optional[str],
        simulation_id: Optional[str] = None,
    ) -> Optional[str]:
        content = msg.get('content', '').strip()
        conv_id = msg.get('conversation_id')
        
        # 0. Segurança: Valida se a conversa pertence ao usuário
        if conv_id:
            def _check_conv():
                return self.db.table('conversations').select('username').eq('id', conv_id).execute().data
            
            conv_rows = await self._execute_db(_check_conv)
            if conv_rows and conv_rows[0]['username'] != username:
                print(f"[ChatSecurity] Alerta: Usuário {username} tentou acessar conversa {conv_id} de {conv_rows[0]['username']}")
                await websocket.send_json({'type': 'error', 'message': 'Acesso negado à conversa.'})
                return pending_drill_target

        is_voice_mode = msg.get('origin') == 'voice'
        msg_type = msg.get('type')

        # 1. Transcrição ou Extração de Arquivo
        if msg_type == 'audio':
            #print(f'[ChatService] Transcrevendo áudio para {username}...')
            audio_bytes = base64.b64decode(msg.get('audio', ''))
            content = await transcribe_audio(audio_bytes, filename='input.webm')
            #print(f'[ChatService] Transcrição concluída: "{content[:50]}..."')
            await websocket.send_json({'type': 'transcription', 'text': content})
        elif msg_type == 'file':
            filename = msg.get('filename', 'file.txt')
            extracted = await self.extract_text_from_file(
                filename, msg.get('content', '')
            )
            caption = msg.get('caption', '').strip()
            content = (
                f'{caption}\n\n[Arquivo: {filename}]\n{extracted}'
                if caption
                else f'[Arquivo: {filename}]\n{extracted}'
            )
            await websocket.send_json(
                {'type': 'status', 'text': f'Arquivo {filename} lido.'}
            )

        if not content:
            return pending_drill_target

        # 2. Drill result
        if pending_drill_target:
            result = match_pronunciation(pending_drill_target, content)
            await websocket.send_json(
                {
                    'type': 'drill_result',
                    'result': result,
                    'target': pending_drill_target,
                }
            )
            pending_drill_target = None

        # 3. Access check
        access = await self.check_access(username)
        if not access['allowed']:
            await websocket.send_json(
                {'type': 'error', 'code': 402, 'detail': 'Limit reached'}
            )
            return pending_drill_target

        # 4. History and background tasks
        history = await load_llm_history(conv_id)
        if not history:
            await auto_title(conv_id, username, content[:50])

        user_audio_b64 = msg.get('audio') if msg_type == 'audio' else None
        await save_message(conv_id, username, 'user', content, audio_b64=user_audio_b64)
        history.append({'role': 'user', 'content': content})

        # Streak e metadados (rápido)
        from app.modules.users.services.streaks import record_study_day
        asyncio.create_task(record_study_day(username))

        # 5. LLM Response - OTIMIZAÇÃO: Busca RAG, Perfil e Sentimento em PARALELO
        # Só buscamos RAG se a mensagem for relevante (mais de 10 caracteres)
        use_rag = len(content) > 10
        
        # Define tarefas para execução paralela
        tasks = [self.get_user_profile(username)]
        
        if use_rag:
            tasks.append(run_in_threadpool(obter_contexto_rag, content))
        else:
            async def _dummy_rag(): return RAGResult(contexto='', fontes='')
            tasks.append(_dummy_rag())

        # Podcasts reais (busca rápida)
        def _fetch_pods():
            return (self.db.table('podcasts').select('title, category').eq('user_id', username).order('created_at', desc=True).limit(5).execute().data)
        tasks.append(self._execute_db(_fetch_pods))

        # Busca prompt de simulação em paralelo se houver
        if simulation_id:
            def _fetch_sim():
                return self.db.table('simulations').select('system_prompt').eq('id', simulation_id).limit(1).execute().data
            tasks.append(self._execute_db(_fetch_sim))
        else:
            async def _dummy_sim(): return []
            tasks.append(_dummy_sim())

        # --- Sprint 5: Análise de Sentimento (PARALELIZADA) ---
        async def _detect_sentiment_task(text: str) -> str:
            from app.modules.chat.services.llm import groq_chat
            prompt = f"Analyze student message and return ONLY ONE word representing their sentiment: FRUSTRATED, EXCITED, CONFUSED, NEUTRAL, or TIRED.\n\nMessage: {text}"
            try:
                # Usa groq_chat diretamente (já é async) e especifica modelo ultra-rápido
                res = await groq_chat([{"role": "user", "content": prompt}], max_tokens=10, temperature=0.1, model='llama-3.1-8b-instant')
                return res.strip().upper()
            except: return "NEUTRAL"
        
        # Truncate content to avoid 413 Payload Too Large when large files are attached
        tasks.append(_detect_sentiment_task(content[:1000]))

        # Executa todas as tarefas de preparação em paralelo
        results = await asyncio.gather(*tasks)
        profile = results[0]
        rag_result = results[1]
        real_podcasts = results[2] or []
        sim_rows = results[3]
        sentiment = results[4]

        sim_prompt = sim_rows[0].get('system_prompt') if sim_rows else None

        sentiment_map = {
            "FRUSTRATED": "Student feels FRUSTRATED. Be extremely patient, empathetic, and encouraging. Use simple words.",
            "EXCITED": "Student is EXCITED! Be very enthusiastic and celebrate their energy.",
            "CONFUSED": "Student is CONFUSED. Explain concepts slowly and clearly. Ask if they need an example.",
            "TIRED": "Student seems TIRED. Keep the conversation light, easy, and supportive.",
            "NEUTRAL": ""
        }
        sentiment_instruction = sentiment_map.get(sentiment, "")

        effective_prompt = build_effective_prompt(
            profile, rag_result.contexto, real_podcasts=real_podcasts
        )
        
        if sentiment_instruction:
            effective_prompt = f"{sentiment_instruction}\n\n{effective_prompt}"

        if sim_prompt:
            # REFORÇO CRÍTICO: Simulação deve ser o foco absoluto e manter nome Tati
            effective_prompt = (
                f"ACT AS THE CHARACTER IN THIS SCENARIO:\n{sim_prompt}\n\n"
                f"STRICT RULES:\n"
                f"1. YOUR NAME IS ALWAYS TATI (or Tatiana). Introduce yourself as Tati.\n"
                f"2. STAY IN CHARACTER at all times.\n"
                f"3. DO NOT discuss anything outside the scenario.\n"
                f"4. Keep responses brief and relevant to the simulation.\n"
                f"5. If the user goes off-topic, politely bring them back to the scenario.\n\n"
                f"--- USER PROFILE ---\n{effective_prompt}"
            )

        if is_voice_mode:
            effective_prompt += '\n\nCRITICAL: Voice mode. Keep responses very short (max 2 sentences).'
            
        effective_prompt += '\n\nCRITICAL: If the user explicitly asks you to generate a PDF or a document, your VERY FIRST characters MUST be EXACTLY the tag "[GENERATE_PDF: filename.pdf]" where filename is related to the topic. DO NOT output any greetings or conversational text before this tag. Everything after this tag will be formatted into a downloadable PDF file. The content of the PDF MUST be entirely in English.'

        await websocket.send_json({'type': 'stream_start', 'conversation_id': conv_id})
        full_response = ''
        # Limite de tokens para velocidade
        max_tokens = 1000 if is_voice_mode else 1500
        
        is_pdf_generation = False
        pdf_filename = f"tati_doc_{conv_id}.pdf"
        buffer = ''
        
        async for token in stream_llm(effective_prompt, history, max_tokens=max_tokens):
            if not is_pdf_generation and len(buffer) < 80:
                buffer += token
                if '[GENERATE_PDF' in buffer and ']' in buffer:
                    is_pdf_generation = True
                    start_idx = buffer.find('[GENERATE_PDF')
                    end_idx = buffer.find(']', start_idx)
                    tag_content = buffer[start_idx:end_idx+1]
                    
                    if ':' in tag_content:
                        pdf_filename = tag_content.split(':', 1)[1].strip(' ]')
                        if not pdf_filename.endswith('.pdf'):
                            pdf_filename += '.pdf'
                            
                    # Get everything after the tag
                    full_response = buffer[end_idx+1:].lstrip()
                elif len(buffer) >= 80 and '[GENERATE_PDF' not in buffer:
                    # Flush buffer as normal text
                    full_response += buffer
                    await websocket.send_json({'type': 'stream_token', 'content': buffer})
            else:
                if is_pdf_generation:
                    full_response += token
                else:
                    full_response += token
                    await websocket.send_json({'type': 'stream_token', 'content': token})
        
        # Se terminou e buffer era menor que 80 e não era PDF
        if not is_pdf_generation and len(buffer) < 80 and buffer:
            if '[GENERATE_PDF' not in buffer:
                full_response += buffer
                await websocket.send_json({'type': 'stream_token', 'content': buffer})

        if is_pdf_generation:
            try:
                from app.shared.services.pdf_generator import generate_report_pdf
                pdf_path = generate_report_pdf(full_response.strip(), filename=pdf_filename)
                with open(pdf_path, 'rb') as f:
                    pdf_b64 = base64.b64encode(f.read()).decode('utf-8')
                await websocket.send_json({
                    'type': 'pdf_generated', 
                    'pdf_b64': pdf_b64, 
                    'filename': pdf_filename,
                    'text': '📄 Document created successfully!'
                })
            except Exception as e:
                print(f"Error generating PDF: {e}")
                await websocket.send_json({'type': 'stream_token', 'content': '\n[Error generating PDF]'})

        await websocket.send_json({'type': 'stream_end'})

        # 6. Finaliza em Background para não travar o socket
        asyncio.create_task(
            self._post_response_tasks(username, content, full_response, websocket, conv_id, is_pdf_generation)
        )

        return pending_drill_target

    async def _run_background_tasks(
        self, username: str, content: str, history: List[Dict[str, Any]]
    ):
        """Executa tarefas pesadas fora do fluxo principal do chat."""
        try:
            # 1. Podcast Discovery
            from app.modules.activities.services.podcast_discovery import discover_personalized_podcasts

            profile = await self.get_user_profile(username)
            await discover_personalized_podcasts(username, username, profile.level)

            # 2. Record Streak
            from app.modules.users.services.streaks import record_study_day
            await record_study_day(username)
        except Exception as e:
            print(f'[Chat Background] Erro: {e}')

    async def _post_response_tasks(
        self, username: str, user_content: str, full_response: str, websocket: WebSocket, conv_id: str, is_pdf_generation: bool = False
    ):
        """Tarefas após o streaming terminar."""
        try:
            # 0. XP (Gamification) - Recompensa por mensagem enviada
            from app.modules.activities.services.gamification_service import GamificationService
            gs = GamificationService()
            asyncio.create_task(gs.award_xp(username, gs.XP_REWARDS['message_sent'], 'Chat activity'))

            # 1. TTS e Salvar (PRIORIDADE)
            audio_b64 = None
            tts_text = ""
            if not is_pdf_generation:
                tts_text = self._clean_tts_text(full_response)
                audio_b64 = await text_to_speech(tts_text)
                
            await save_message(
                conv_id, username, 'assistant', full_response, audio_b64=audio_b64
            )

            # Envia áudio imediatamente se for modo voz (opcional, já vai via WS se for o caso)
            if audio_b64 and not is_pdf_generation:
                await websocket.send_json({'type': 'audio_response', 'audio': audio_b64, 'content': full_response})

            # 1. Podcast Discovery (Atraso de 2s para não impactar a resposta imediata)
            async def _delayed_discovery():
                await asyncio.sleep(2)
                from app.modules.activities.services.podcast_discovery import discover_personalized_podcasts
                profile = await self.get_user_profile(username)
                await discover_personalized_podcasts(username, username, profile.level)
            
            asyncio.create_task(_delayed_discovery())

            # 2. Trophies
            from app.modules.activities.services.trophy_service import check_chat_trophies

            await run_in_threadpool(check_chat_trophies, username)

            # 3. AutoExercise, Error Logging e Weekly Plan Tracking
            from app.modules.activities.services.error_log_service import error_log_service
            from app.modules.chat.services.semantic_judge import semantic_judge
            
            # Executa tarefas de background de forma assíncrona
            # AGORA PASSANDO O CONTEÚDO DO USUÁRIO CORRETAMENTE
            asyncio.create_task(error_log_service.extract_and_log_errors(username, user_content, full_response))
            asyncio.create_task(semantic_judge.check_topics_completion(username, f"User: {user_content}\nTati: {full_response}"))

        except Exception as e:
            print(f'[Chat Post-Response] Erro: {e}')


    async def extract_text_from_file(self, filename: str, content_b64: str) -> str:
        ext = filename.lower().rsplit('.', 1)[-1] if '.' in filename else ''
        
        # Apenas permite extração de texto de arquivos suportados
        text_exts = {'txt', 'md', 'csv', 'json', 'py', 'js', 'html', 'css'}
        
        image_exts = {'png', 'jpg', 'jpeg', 'gif', 'bmp', 'webp', 'svg'}
        if ext in image_exts:
            return f'[Imagem anexada: {filename}. A Teacher Tati ainda não possui visão computacional, mas sabe que você enviou esta imagem.]'
            
        try:
            file_bytes = base64.b64decode(content_b64)
            if ext == 'pdf':
                reader = pypdf.PdfReader(io.BytesIO(file_bytes))
                return '\n'.join(
                    p.extract_text() for p in reader.pages if p.extract_text()
                ).replace('\x00', '')
            if ext == 'docx':
                doc = docx.Document(io.BytesIO(file_bytes))
                return '\n'.join(p.text for p in doc.paragraphs).replace('\x00', '')
            if ext == 'pptx':
                from pptx import Presentation
                prs = Presentation(io.BytesIO(file_bytes))
                text_runs = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_runs.append(shape.text)
                return '\n'.join(text_runs).replace('\x00', '')
            if ext == 'xlsx':
                import openpyxl
                wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
                lines = []
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        row_text = ' | '.join([str(cell) for cell in row if cell is not None])
                        if row_text:
                            lines.append(row_text)
                return '\n'.join(lines).replace('\x00', '')
                
            if ext in text_exts:
                text_content = file_bytes.decode('utf-8', errors='ignore').replace('\x00', '')
                if len(text_content) > 50000:
                    return text_content[:50000] + '\n\n[Texto truncado (muito grande)]'
                return text_content
                
            return f'[Arquivo anexado: {filename}. O conteúdo deste formato não pode ser lido diretamente pelo chat no momento.]'
        except Exception as e:
            return f'[Erro ao ler arquivo: {e}]'

    def _clean_tts_text(self, text: str) -> str:
        text = text.replace('*', '').replace('#', '').replace('_', '')
        text = re.sub(r'\[DRILL:.*?\]', '', text)
        return text.strip() or 'Please, repeat.'

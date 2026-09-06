import base64
import io
import json
import logging
import os
import re
import shutil
import subprocess
import uuid
from typing import Any, Dict, List, Optional, Tuple

from django.conf import settings

from .audio_service import strip_emojis

logger = logging.getLogger(__name__)

def get_media_docs_dir() -> str:
    media_root = "media"
    try:
        if settings.configured and hasattr(settings, "MEDIA_ROOT"):
            media_root = str(settings.MEDIA_ROOT)
    except Exception:
        pass
    path = os.path.join(media_root, "generated_docs")
    os.makedirs(path, exist_ok=True)
    return path


class DocumentService:
    """
    Serviço completo para:
    1. Leitura e extração integral de arquivos enviados no chat (máx 3): PDF, Word (.docx, .doc),
       PowerPoint (.pptx, .ppt), imagens (visão) e arquivos de texto.
    2. Geração formatada assíncrona de arquivos solicitados pelo usuário nos formatos:
       - PDF (.pdf) como padrão (se não informado)
       - Word (.docx / .doc)
       - PowerPoint (.pptx / .ppt)
    3. Suporte a pré-visualização no navegador e download direto, persistindo para que o usuário
       possa sair do chat e ao voltar o arquivo esteja pronto para abrir no navegador.
    """

    # ─────────────────────────────────────────────────────────────────────────
    # 1. LEITURA INTEGRAL DE ARQUIVOS (MÁXIMO 3)
    # ─────────────────────────────────────────────────────────────────────────

    @classmethod
    def read_uploaded_files(cls, files: List[Dict[str, Any]]) -> str:
        """
        Recebe até 3 arquivos enviados pelo aluno e extrai o conteúdo textual integral de cada um.
        """
        if not files:
            return ""

        extracted_texts = []
        # Limita estritamente a 3 arquivos
        safe_files = files[:3]

        for idx, f in enumerate(safe_files, start=1):
            filename = f.get("filename") or f.get("name") or f"arquivo_{idx}"
            b64_data = f.get("base64") or f.get("file") or f.get("content") or ""
            
            # Limpa prefixos data URI (data:application/pdf;base64,...)
            if "," in b64_data:
                b64_data = b64_data.split(",", 1)[1]

            if not b64_data:
                continue

            try:
                raw_bytes = base64.b64decode(b64_data)
                ext = os.path.splitext(filename)[1].lower()
                text = cls._extract_text_from_bytes(filename, ext, raw_bytes)
                if text and text.strip():
                    extracted_texts.append(
                        f"=== [ARQUIVO {idx}/3: {filename}] ===\n{text.strip()}\n"
                    )
                else:
                    extracted_texts.append(
                        f"=== [ARQUIVO {idx}/3: {filename}] ===\n(Arquivo processado, sem texto legível extraído)\n"
                    )
            except Exception as e:
                logger.error(f"[DocReader] Erro ao processar '{filename}': {e}")
                extracted_texts.append(
                    f"=== [ARQUIVO {idx}/3: {filename}] ===\n(Erro ao extrair conteúdo do arquivo: {str(e)})\n"
                )

        if not extracted_texts:
            return ""

        return (
            "\n\n=== CONTEÚDO INTEGRAL DOS DOCUMENTOS ENVIADOS PELO ALUNO (MÁXIMO 3) ===\n"
            + "\n".join(extracted_texts)
            + "\n=========================================================================\n"
        )

    @classmethod
    def _extract_text_from_bytes(cls, filename: str, ext: str, raw_bytes: bytes) -> str:
        """Extrai texto dependendo da extensão do arquivo."""
        # PDF
        if ext == ".pdf":
            return cls._extract_from_pdf(raw_bytes)

        # Word DOCX
        if ext == ".docx":
            return cls._extract_from_docx(raw_bytes)

        # Word DOC (legado)
        if ext == ".doc":
            return cls._extract_from_doc_legacy(filename, raw_bytes)

        # PowerPoint PPTX
        if ext == ".pptx":
            return cls._extract_from_pptx(raw_bytes)

        # PowerPoint PPT (legado)
        if ext == ".ppt":
            return cls._extract_from_ppt_legacy(filename, raw_bytes)

        # Arquivos de texto / código
        if ext in [".txt", ".md", ".csv", ".json", ".xml", ".html", ".py", ".js", ".ts", ".srt"]:
            return cls._extract_from_text(raw_bytes)

        # Imagens (OCR / Visão com Gemini)
        if ext in [".png", ".jpg", ".jpeg", ".webp", ".bmp"]:
            return cls._extract_from_image(ext, raw_bytes)

        # Fallback genérico para texto utf-8/latin1
        return cls._extract_from_text(raw_bytes)

    @staticmethod
    def _extract_from_pdf(raw_bytes: bytes) -> str:
        pages_text = []
        try:
            import pypdf

            reader = pypdf.PdfReader(io.BytesIO(raw_bytes))
            for i, page in enumerate(reader.pages):
                txt = page.extract_text() or ""
                if txt.strip():
                    pages_text.append(f"[Página {i+1}]\n{txt.strip()}")
        except Exception as e:
            logger.warning(f"[DocReader] pypdf falhou: {e}")

        # Se pypdf falhou ou extraiu pouco, tenta PyMuPDF (fitz)
        if not pages_text or len("\n".join(pages_text)) < 50:
            try:
                import fitz

                doc = fitz.open(stream=raw_bytes, filetype="pdf")
                pages_text = []
                for i in range(len(doc)):
                    txt = doc[i].get_text() or ""
                    if txt.strip():
                        pages_text.append(f"[Página {i+1}]\n{txt.strip()}")
            except Exception as e:
                logger.warning(f"[DocReader] PyMuPDF falhou: {e}")

        return "\n\n".join(pages_text)

    @staticmethod
    def _extract_from_docx(raw_bytes: bytes) -> str:
        try:
            import docx

            doc = docx.Document(io.BytesIO(raw_bytes))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        paragraphs.append(row_text)
            return "\n".join(paragraphs)
        except Exception as e:
            logger.error(f"[DocReader] Erro em docx: {e}")
            return ""

    @classmethod
    def _extract_from_doc_legacy(cls, filename: str, raw_bytes: bytes) -> str:
        # Tenta converter via LibreOffice headless se disponível
        soffice = cls._get_soffice_cmd()
        if soffice:
            import tempfile

            with tempfile.TemporaryDirectory() as tmpdir:
                in_path = os.path.join(tmpdir, filename)
                with open(in_path, "wb") as f:
                    f.write(raw_bytes)
                cmd = [soffice, "--headless", "--convert-to", "docx", in_path, "--outdir", tmpdir]
                res = subprocess.run(cmd, capture_output=True, timeout=30)
                out_name = os.path.splitext(filename)[0] + ".docx"
                out_path = os.path.join(tmpdir, out_name)
                if os.path.exists(out_path):
                    with open(out_path, "rb") as f:
                        return cls._extract_from_docx(f.read())

        # Fallback para extração de strings legíveis
        return cls._extract_printable_strings(raw_bytes)

    @staticmethod
    def _extract_from_pptx(raw_bytes: bytes) -> str:
        try:
            import pptx

            prs = pptx.Presentation(io.BytesIO(raw_bytes))
            slides_text = []
            for i, slide in enumerate(prs.slides, start=1):
                slide_lines = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            line = "".join(run.text for run in paragraph.runs).strip()
                            if line:
                                slide_lines.append(line)
                if slide_lines:
                    slides_text.append(f"[Slide {i}]\n" + "\n".join(slide_lines))
            return "\n\n".join(slides_text)
        except Exception as e:
            logger.error(f"[DocReader] Erro em pptx: {e}")
            return ""

    @classmethod
    def _extract_from_ppt_legacy(cls, filename: str, raw_bytes: bytes) -> str:
        soffice = cls._get_soffice_cmd()
        if soffice:
            import tempfile

            with tempfile.TemporaryDirectory() as tmpdir:
                in_path = os.path.join(tmpdir, filename)
                with open(in_path, "wb") as f:
                    f.write(raw_bytes)
                cmd = [soffice, "--headless", "--convert-to", "pptx", in_path, "--outdir", tmpdir]
                subprocess.run(cmd, capture_output=True, timeout=30)
                out_name = os.path.splitext(filename)[0] + ".pptx"
                out_path = os.path.join(tmpdir, out_name)
                if os.path.exists(out_path):
                    with open(out_path, "rb") as f:
                        return cls._extract_from_pptx(f.read())

        return cls._extract_printable_strings(raw_bytes)

    @staticmethod
    def _extract_from_text(raw_bytes: bytes) -> str:
        for enc in ["utf-8", "latin-1", "cp1252"]:
            try:
                return raw_bytes.decode(enc)
            except UnicodeDecodeError:
                pass
        return raw_bytes.decode("utf-8", errors="replace")

    @staticmethod
    def _extract_from_image(ext: str, raw_bytes: bytes) -> str:
        gemini_key = os.getenv("GEMINI_API_KEY") or os.getenv("GEMINI_API_KEY_1")
        if not gemini_key:
            return "(Imagem recebida para análise visual)"

        try:
            import google.generativeai as genai

            genai.configure(api_key=gemini_key)
            model = genai.GenerativeModel("gemini-1.5-flash")
            mime_map = {
                ".png": "image/png",
                ".jpg": "image/jpeg",
                ".jpeg": "image/jpeg",
                ".webp": "image/webp",
                ".bmp": "image/bmp",
            }
            mime = mime_map.get(ext, "image/jpeg")
            response = model.generate_content(
                [
                    {"mime_type": mime, "data": raw_bytes},
                    "Transcreva e descreva com detalhes pedagógicos todo o texto e informações contidas nesta imagem para Teacher Tatiana Duarte e seu aluno.",
                ]
            )
            return response.text or ""
        except Exception as e:
            logger.warning(f"[DocReader] Gemini Vision falhou: {e}")
            return "(Imagem recebida pelo aluno)"

    @staticmethod
    def _extract_printable_strings(raw_bytes: bytes) -> str:
        try:
            text = "".join(
                chr(b) if (32 <= b <= 126 or b in [10, 13]) else " " for b in raw_bytes
            )
            words = [w for w in text.split() if len(w) > 3]
            return " ".join(words[:2000])
        except Exception:
            return ""

    @staticmethod
    def _get_soffice_cmd() -> Optional[str]:
        soffice_path = (
            "libreoffice"
            if shutil.which("libreoffice")
            else ("soffice" if shutil.which("soffice") else None)
        )
        if not soffice_path:
            for p in [
                "/usr/bin/libreoffice",
                "/usr/bin/soffice",
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            ]:
                if os.path.exists(p):
                    return p
        return soffice_path

    # ─────────────────────────────────────────────────────────────────────────
    # 2. DETECÇÃO DE FORMATO SOLICITADO & SE DEVE GERAR DOCUMENTO
    # ─────────────────────────────────────────────────────────────────────────

    @classmethod
    def detect_target_format(cls, user_text: str) -> str:
        """
        Detecta se o usuário pediu pdf, doc, word, ppt ou pptx.
        Regra solicitada: 'no formato que ela pedir dentre esses, se não informar, será pdf como padrão'.
        """
        lower = (user_text or "").lower()

        # PowerPoint
        if any(k in lower for k in ["ppt", "pptx", "powerpoint", "slides", "apresentação", "apresentacao"]):
            return "pptx"

        # Word / Doc
        if any(k in lower for k in ["word", "doc", "docx", "documento word"]):
            return "docx"

        # Padrão: PDF
        return "pdf"

    @classmethod
    def should_generate_document(cls, user_text: str, num_files: int = 0) -> Tuple[bool, str]:
        """
        Determina se uma solicitação deve gerar um arquivo formatado para download e abertura no navegador.
        Retorna (deve_gerar: bool, formato_alvo: str).
        """
        target_format = cls.detect_target_format(user_text)
        lower = (user_text or "").lower()

        keywords = [
            "crie um arquivo", "criar arquivo", "gerar arquivo", "gere um arquivo",
            "crie um documento", "gerar documento", "gere documento", "criar documento",
            "faça um resumo", "fazer resumo", "gerar resumo", "crie um resumo",
            "crie slides", "gerar slides", "fazer slides", "apresentação", "apresentacao",
            "em pdf", "em word", "em doc", "em docx", "em ppt", "em pptx",
            "no formato", "salve em", "exportar", "baixe", "download",
            "lista de exercícios", "exercicios", "apostila", "relatório", "relatorio",
            "formate como", "gere um pdf", "crie um pdf", "gerar pdf", "criar pdf",
            "arquivo formatado", "formate o arquivo", "leia e crie", "leia e gere"
        ]

        if any(k in lower for k in keywords):
            return True, target_format

        # Se enviou arquivos e escreveu instrução no chat (ex: 'analise e estruture', 'resuma')
        if num_files > 0 and len(user_text.strip()) > 5:
            # Qualquer instrução com arquivos enviados gera o documento consolidado
            return True, target_format

        return False, target_format

    # ─────────────────────────────────────────────────────────────────────────
    # 3. GERAÇÃO FORMATADA DO DOCUMENTO (PDF, DOCX, PPTX)
    # ─────────────────────────────────────────────────────────────────────────

    @classmethod
    def generate_document_from_instruction(
        cls,
        user_text: str,
        files_extracted_text: str,
        student_name: str,
        target_format: str = "pdf",
    ) -> Dict[str, Any]:
        """
        Cria de forma estruturada e assíncrona o documento formatado solicitado.
        Retorna dicionário com metadados, URLs públicas, nome de arquivo e status.
        """
        # 1. Estrutura o conteúdo usando Groq/Gemini
        structured_data = cls._generate_structured_content(
            user_text=user_text,
            files_extracted_text=files_extracted_text,
            student_name=student_name,
            target_format=target_format,
        )

        doc_title = structured_data.get("title") or "Material de Estudo Teacher Tati"
        safe_title = re.sub(r"[^a-zA-Z0-9_\-]", "_", doc_title)[:35].strip("_") or "Teacher_Tati_Document"
        doc_uuid = uuid.uuid4().hex[:10]
        base_filename = f"{safe_title}_{doc_uuid}"

        # 2. Gera o arquivo de acordo com o formato solicitado
        file_path = ""
        actual_format = target_format.lower()
        docs_dir = get_media_docs_dir()

        if actual_format == "docx" or actual_format == "doc":
            actual_format = "docx"
            file_path = os.path.join(docs_dir, f"{base_filename}.docx")
            cls._build_docx(structured_data, file_path, student_name)
        elif actual_format == "pptx" or actual_format == "ppt":
            actual_format = "pptx"
            file_path = os.path.join(docs_dir, f"{base_filename}.pptx")
            cls._build_pptx(structured_data, file_path, student_name)
        else:
            actual_format = "pdf"
            file_path = os.path.join(docs_dir, f"{base_filename}.pdf")
            cls._build_pdf(structured_data, file_path, student_name)

        # 3. Calcula tamanho legível
        file_size_bytes = os.path.getsize(file_path) if os.path.exists(file_path) else 0
        size_str = cls._format_file_size(file_size_bytes)

        # 4. URLs de acesso no navegador
        rel_filename = os.path.basename(file_path)
        media_url = f"/media/generated_docs/{rel_filename}"

        # Se for DOCX ou PPTX, tenta gerar um PDF companheiro via LibreOffice para abrir direto na aba
        preview_url = media_url
        if actual_format in ["docx", "pptx"]:
            companion_pdf = cls._try_convert_to_pdf(file_path)
            if companion_pdf:
                preview_url = f"/media/generated_docs/{os.path.basename(companion_pdf)}"

        # 5. Base64 para download instantâneo se for PDF (ou DOCX/PPTX convertido)
        pdf_b64 = ""
        if actual_format == "pdf" and os.path.exists(file_path):
            with open(file_path, "rb") as pf:
                pdf_b64 = base64.b64encode(pf.read()).decode("utf-8")
        elif actual_format in ["docx", "pptx"] and "companion_pdf" in locals() and companion_pdf and os.path.exists(companion_pdf):
            with open(companion_pdf, "rb") as pf:
                pdf_b64 = base64.b64encode(pf.read()).decode("utf-8")

        return {
            "id": doc_uuid,
            "title": doc_title,
            "filename": rel_filename,
            "format": actual_format,
            "url": media_url,
            "preview_url": preview_url,
            "size": size_str,
            "pdf_b64": pdf_b64,
            "summary": structured_data.get("summary", ""),
        }

    @classmethod
    def _generate_structured_content(
        cls,
        user_text: str,
        files_extracted_text: str,
        student_name: str,
        target_format: str,
    ) -> Dict[str, Any]:
        """Usa IA para criar tópicos, parágrafos e conteúdo pedagógico enriquecido."""
        safe_files_text = files_extracted_text[:12000] if files_extracted_text else ""
        prompt = (
            f"You are Teacher Tatiana Duarte's pedagogical content assistant.\\n"
            f"The student '{student_name}' requested a formatted document in format: '{target_format.upper()}'.\\n"
            f"Student's instruction: \"{user_text}\"\\n\\n"
            f"{safe_files_text}\\n\\n"
            f"Based on the files (if provided) and the student's instructions, create a comprehensive, highly organized educational document.\\n"
            f"Return a strict JSON object with:\\n"
            f"- 'title': An inspiring, professional title for the document\\n"
            f"- 'subtitle': A descriptive subtitle\\n"
            f"- 'summary': A brief introductory overview (2-3 sentences)\\n"
            f"- 'sections': A list of 3 to 6 sections, each having:\\n"
            f"    - 'heading': Section title\\n"
            f"    - 'paragraphs': List of 1-3 well-written text paragraphs in English (or Portuguese if requested)\\n"
            f"    - 'bullet_points': List of 2-5 clear bullet points (key rules, examples, takeaways or exercises)\\n"
            f"- 'key_takeaways': List of 3-4 bullet takeaways or practice tips at the end\\n"
            f"Strict rules: Valid JSON ONLY. No markdown code formatting ticks. No emojis."
        )

        from .word_service import WordLookupService

        keys = WordLookupService._get_groq_keys()

        for key in keys:
            try:
                from groq import Groq

                client = Groq(api_key=key, timeout=14.0)
                completion = client.chat.completions.create(
                    model="openai/gpt-oss-20b",
                    messages=[
                        {
                            "role": "system",
                            "content": "You are a professional educational material formatter. Output strictly valid JSON without any emojis.",
                        },
                        {"role": "user", "content": prompt},
                    ],
                    response_format={"type": "json_object"},
                    temperature=0.3,
                    max_tokens=2500,
                )
                raw_json = completion.choices[0].message.content or "{}"
                data = json.loads(raw_json)
                if data.get("title") and data.get("sections"):
                    return data
            except Exception as e:
                logger.warning(f"[DocGen] Groq key failed for structured content: {e}")

        # Fallback estruturado básico caso offline
        return {
            "title": f"Study Guide & Summary — {student_name}",
            "subtitle": "Created by Teacher Tatiana Duarte AI",
            "summary": "This document consolidates your requested materials, key concepts, and practical English exercises.",
            "sections": [
                {
                    "heading": "1. Overview and Core Concepts",
                    "paragraphs": ["Review the core structures discussed in this practice session."],
                    "bullet_points": [
                        "Focus on natural everyday conversational patterns.",
                        "Pay attention to verb tenses and practical idioms.",
                    ],
                },
                {
                    "heading": "2. Practical Application & Examples",
                    "paragraphs": ["Put the concepts into practice with consistent active recall."],
                    "bullet_points": [
                        "Practice reading aloud to improve pronunciation and confidence.",
                        "Formulate sentences reflecting your real-life daily routine.",
                    ],
                },
            ],
            "key_takeaways": [
                "Consistency beats intensity in language acquisition.",
                "Review these notes regularly before our next chat session.",
            ],
        }

    # ─────────────────────────────────────────────────────────────────────────
    # 4. CONSTRUTORES DE ARQUIVO (PDF, DOCX, PPTX)
    # ─────────────────────────────────────────────────────────────────────────

    @classmethod
    def _build_pdf(cls, data: Dict[str, Any], output_path: str, student_name: str):
        """Gera PDF de alta qualidade com ReportLab."""
        from reportlab.lib.colors import HexColor
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.platypus import HRFlowable, KeepTogether, Paragraph, SimpleDocTemplate, Spacer

        doc = SimpleDocTemplate(
            output_path,
            pagesize=letter,
            rightMargin=45,
            leftMargin=45,
            topMargin=45,
            bottomMargin=45,
        )

        styles = getSampleStyleSheet()
        primary_color = HexColor("#4F46E5")
        dark_color = HexColor("#0F172A")
        slate_color = HexColor("#334155")
        light_bg = HexColor("#F8FAFC")

        title_style = ParagraphStyle(
            "DocTitle",
            parent=styles["Normal"],
            fontName="Helvetica-Bold",
            fontSize=22,
            leading=26,
            textColor=primary_color,
            spaceAfter=6,
        )
        subtitle_style = ParagraphStyle(
            "DocSubTitle",
            parent=styles["Normal"],
            fontName="Helvetica",
            fontSize=11,
            leading=15,
            textColor=slate_color,
            spaceAfter=14,
        )
        meta_style = ParagraphStyle(
            "DocMeta",
            parent=styles["Normal"],
            fontName="Helvetica-Bold",
            fontSize=8,
            leading=11,
            textColor=primary_color,
            textTransform="uppercase",
            spaceAfter=4,
        )
        h1_style = ParagraphStyle(
            "DocH1",
            parent=styles["Normal"],
            fontName="Helvetica-Bold",
            fontSize=13,
            leading=17,
            textColor=dark_color,
            spaceBefore=12,
            spaceAfter=6,
        )
        body_style = ParagraphStyle(
            "DocBody",
            parent=styles["Normal"],
            fontName="Helvetica",
            fontSize=9.5,
            leading=14,
            textColor=slate_color,
            spaceAfter=6,
        )
        bullet_style = ParagraphStyle(
            "DocBullet",
            parent=styles["Normal"],
            fontName="Helvetica",
            fontSize=9,
            leading=13,
            textColor=dark_color,
            leftIndent=15,
            firstLineIndent=-10,
            spaceAfter=4,
        )

        story = []

        # Cabeçalho da Teacher Tati
        story.append(Paragraph("TEACHER TATI AI • PEDAGOGICAL MATERIAL", meta_style))
        title = strip_emojis(data.get("title") or "Material de Estudo")
        story.append(Paragraph(title, title_style))
        subtitle = strip_emojis(data.get("subtitle") or f"Preparado especialmente para {student_name}")
        story.append(Paragraph(subtitle, subtitle_style))
        story.append(HRFlowable(width="100%", thickness=1.5, color=primary_color, spaceAfter=14))

        # Resumo / Summary
        summary = strip_emojis(data.get("summary") or "")
        if summary:
            story.append(Paragraph(f"<b>Overview:</b> {summary}", body_style))
            story.append(Spacer(1, 8))

        # Seções
        for sec in data.get("sections", []):
            sec_elements = []
            heading = strip_emojis(sec.get("heading") or "")
            if heading:
                sec_elements.append(Paragraph(heading, h1_style))

            for p in sec.get("paragraphs", []):
                clean_p = strip_emojis(p)
                if clean_p:
                    sec_elements.append(Paragraph(clean_p, body_style))

            for b in sec.get("bullet_points", []):
                clean_b = strip_emojis(b)
                if clean_b:
                    sec_elements.append(Paragraph(f"• {clean_b}", bullet_style))

            sec_elements.append(Spacer(1, 10))
            story.append(KeepTogether(sec_elements))

        # Key Takeaways
        takeaways = data.get("key_takeaways", [])
        if takeaways:
            story.append(Spacer(1, 6))
            story.append(HRFlowable(width="100%", thickness=1, color=HexColor("#CBD5E1"), spaceAfter=10))
            story.append(Paragraph("Key Takeaways & Action Points", h1_style))
            for tk in takeaways:
                clean_tk = strip_emojis(tk)
                if clean_tk:
                    story.append(Paragraph(f"✔ {clean_tk}", bullet_style))

        doc.build(story)

    @classmethod
    def _build_docx(cls, data: Dict[str, Any], output_path: str, student_name: str):
        """Gera arquivo Microsoft Word (.docx) formatado."""
        import docx
        from docx.shared import Inches, Pt, RGBColor

        doc = docx.Document()

        # Margens
        for section in doc.sections:
            section.top_margin = Inches(0.8)
            section.bottom_margin = Inches(0.8)
            section.left_margin = Inches(0.8)
            section.right_margin = Inches(0.8)

        primary_rgb = RGBColor(79, 70, 229)
        slate_rgb = RGBColor(71, 85, 105)

        # Header da instituição
        header_para = doc.add_paragraph()
        r_head = header_para.add_run("TEACHER TATI AI • PEDAGOGICAL MATERIAL")
        r_head.font.size = Pt(8.5)
        r_head.font.bold = True
        r_head.font.color.rgb = primary_rgb

        # Título
        title_para = doc.add_heading(level=1)
        r_title = title_para.add_run(strip_emojis(data.get("title") or "Material de Estudo"))
        r_title.font.size = Pt(20)
        r_title.font.bold = True
        r_title.font.color.rgb = primary_rgb

        # Subtítulo
        sub_para = doc.add_paragraph()
        r_sub = sub_para.add_run(strip_emojis(data.get("subtitle") or f"Preparado para {student_name}"))
        r_sub.font.size = Pt(11)
        r_sub.font.italic = True
        r_sub.font.color.rgb = slate_rgb

        # Resumo
        summary = strip_emojis(data.get("summary") or "")
        if summary:
            sum_para = doc.add_paragraph()
            r_bold = sum_para.add_run("Overview: ")
            r_bold.font.bold = True
            sum_para.add_run(summary)

        # Seções
        for sec in data.get("sections", []):
            heading = strip_emojis(sec.get("heading") or "")
            if heading:
                h = doc.add_heading(heading, level=2)
                for run in h.runs:
                    run.font.color.rgb = RGBColor(15, 23, 42)

            for p in sec.get("paragraphs", []):
                clean_p = strip_emojis(p)
                if clean_p:
                    doc.add_paragraph(clean_p)

            for b in sec.get("bullet_points", []):
                clean_b = strip_emojis(b)
                if clean_b:
                    doc.add_paragraph(clean_b, style="List Bullet")

        # Conclusão / Key Takeaways
        takeaways = data.get("key_takeaways", [])
        if takeaways:
            h_end = doc.add_heading("Key Takeaways & Action Points", level=2)
            for run in h_end.runs:
                run.font.color.rgb = primary_rgb
            for tk in takeaways:
                clean_tk = strip_emojis(tk)
                if clean_tk:
                    doc.add_paragraph(clean_tk, style="List Bullet")

        doc.save(output_path)

    @classmethod
    def _build_pptx(cls, data: Dict[str, Any], output_path: str, student_name: str):
        """Gera apresentação de slides PowerPoint (.pptx) formatada."""
        import pptx
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        prs = pptx.Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]
        primary_color = RGBColor(79, 70, 229)
        dark_color = RGBColor(15, 23, 42)
        white_color = RGBColor(255, 255, 255)
        slate_color = RGBColor(71, 85, 105)

        # SLIDE 1: Capa (Título)
        slide1 = prs.slides.add_slide(blank_layout)
        # Fundo elegante com shape escuro
        bg_shape = slide1.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = dark_color
        bg_shape.line.color.rgb = dark_color

        tx_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11), Inches(3.2))
        tf = tx_box.text_frame
        tf.word_wrap = True

        p_pre = tf.paragraphs[0]
        p_pre.text = "TEACHER TATI AI • PEDAGOGICAL PRESENTATION"
        p_pre.font.size = Pt(12)
        p_pre.font.bold = True
        p_pre.font.color.rgb = RGBColor(129, 140, 248)
        p_pre.space_after = Pt(14)

        p_title = tf.add_paragraph()
        p_title.text = strip_emojis(data.get("title") or "Mastering English")
        p_title.font.size = Pt(36)
        p_title.font.bold = True
        p_title.font.color.rgb = white_color
        p_title.space_after = Pt(12)

        p_sub = tf.add_paragraph()
        p_sub.text = strip_emojis(data.get("subtitle") or f"Interactive Guide prepared for {student_name}")
        p_sub.font.size = Pt(16)
        p_sub.font.color.rgb = RGBColor(203, 213, 225)

        # SLIDES DE CONTEÚDO (1 slide por seção)
        for idx, sec in enumerate(data.get("sections", []), start=1):
            slide = prs.slides.add_slide(blank_layout)

            # Barra superior
            top_bar = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(0.15))
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = primary_color
            top_bar.line.color.rgb = primary_color

            # Título da Seção
            h_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.3), Inches(1.0))
            h_tf = h_box.text_frame
            h_tf.word_wrap = True
            h_p = h_tf.paragraphs[0]
            h_p.text = strip_emojis(sec.get("heading") or f"Topic {idx}")
            h_p.font.size = Pt(24)
            h_p.font.bold = True
            h_p.font.color.rgb = dark_color

            # Conteúdo em Parágrafos e Bullets
            c_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.8))
            c_tf = c_box.text_frame
            c_tf.word_wrap = True

            first = True
            for p in sec.get("paragraphs", []):
                clean_p = strip_emojis(p)
                if clean_p:
                    para = c_tf.paragraphs[0] if first else c_tf.add_paragraph()
                    first = False
                    para.text = clean_p
                    para.font.size = Pt(15)
                    para.font.color.rgb = slate_color
                    para.space_after = Pt(14)

            for b in sec.get("bullet_points", []):
                clean_b = strip_emojis(b)
                if clean_b:
                    para = c_tf.paragraphs[0] if first else c_tf.add_paragraph()
                    first = False
                    para.text = f"• {clean_b}"
                    para.font.size = Pt(14)
                    para.font.color.rgb = dark_color
                    para.space_after = Pt(8)

        # SLIDE FINAL: Key Takeaways
        takeaways = data.get("key_takeaways", [])
        if takeaways:
            slide_end = prs.slides.add_slide(blank_layout)
            h_box = slide_end.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.3), Inches(1.0))
            h_tf = h_box.text_frame
            h_p = h_tf.paragraphs[0]
            h_p.text = "Key Takeaways & Summary"
            h_p.font.size = Pt(28)
            h_p.font.bold = True
            h_p.font.color.rgb = primary_color

            c_box = slide_end.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(4.5))
            c_tf = c_box.text_frame
            c_tf.word_wrap = True
            for i, tk in enumerate(takeaways):
                clean_tk = strip_emojis(tk)
                if clean_tk:
                    para = c_tf.paragraphs[0] if i == 0 else c_tf.add_paragraph()
                    para.text = f"✔  {clean_tk}"
                    para.font.size = Pt(16)
                    para.font.bold = True
                    para.font.color.rgb = dark_color
                    para.space_after = Pt(16)

        prs.save(output_path)

    @classmethod
    def _try_convert_to_pdf(cls, file_path: str) -> Optional[str]:
        """Tenta gerar PDF companheiro para preview no navegador usando LibreOffice."""
        soffice = cls._get_soffice_cmd()
        if not soffice or not os.path.exists(file_path):
            return None

        try:
            out_dir = os.path.dirname(file_path)
            cmd = [
                soffice,
                "--headless",
                "--invisible",
                "--convert-to",
                "pdf",
                file_path,
                "--outdir",
                out_dir,
            ]
            res = subprocess.run(cmd, capture_output=True, timeout=30)
            base_name = os.path.splitext(os.path.basename(file_path))[0]
            pdf_path = os.path.join(out_dir, f"{base_name}.pdf")
            if os.path.exists(pdf_path):
                return pdf_path
        except Exception as e:
            logger.warning(f"[DocConverter] Conversão para PDF companheiro falhou: {e}")
        return None

    @staticmethod
    def _format_file_size(num_bytes: int) -> str:
        if num_bytes < 1024:
            return f"{num_bytes} B"
        elif num_bytes < 1024 * 1024:
            return f"{num_bytes / 1024:.1f} KB"
        else:
            return f"{num_bytes / (1024 * 1024):.1f} MB"
